from __future__ import annotations

import json
from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
ASSETS_DIR = OUT_DIR / "assets"
PPTX_PATH = OUT_DIR / "interview-deck.pptx"
VERIFY_PATH = OUT_DIR / "verification.json"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

# ── Palette
BLUE = RGBColor(0x0D, 0x9A, 0xBF)
BLUE_DARK = RGBColor(0x07, 0x64, 0x7E)
SLATE = RGBColor(0x2B, 0x37, 0x48)
SLATE_LIGHT = RGBColor(0x64, 0x74, 0x8B)
BG = RGBColor(0xED, 0xF2, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDD, 0xE5, 0xED)
TINT = RGBColor(0xE6, 0xF4, 0xF8)
TINT_BORDER = RGBColor(0xA9, 0xD7, 0xE5)
GREEN = RGBColor(0x15, 0x8B, 0x62)
SKY = RGBColor(0xCF, 0xEC, 0xF4)

QR_URLS = {
    "evidence": "https://lucasung-debug.github.io/hermes-ops-dashboard-page/",
    "portfolio": "https://smjportfolio.com/",
}

RECT = MSO_AUTO_SHAPE_TYPE.RECTANGLE
RRECT = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE


def font_name() -> str:
    fonts_dir = Path("C:/Windows/Fonts")
    if fonts_dir.exists() and list(fonts_dir.glob("Pretendard*")):
        return "Pretendard"
    return "Malgun Gothic"


FONT = font_name()


# ── low-level helpers ────────────────────────────────────────────────
def _sub(parent, tag, **attrs):
    el = parent.makeelement(qn(tag), {k: str(v) for k, v in attrs.items()})
    parent.append(el)
    return el


def add_shadow(shape, blur=90000, dist=34000, color="8FA3B8", alpha=42000):
    """Soft drop shadow via raw DrawingML."""
    spPr = shape._element.spPr
    for e in spPr.findall(qn("a:effectLst")):
        spPr.remove(e)
    eff = _sub(spPr, "a:effectLst")
    shd = _sub(eff, "a:outerShdw", blurRad=blur, dist=dist, dir=5400000, rotWithShape=0)
    clr = _sub(shd, "a:srgbClr", val=color)
    _sub(clr, "a:alpha", val=alpha)


def add_text(slide, text, x, y, w, h, size, color=SLATE, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=None,
             letter=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    frame.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=13, color=SLATE_LIGHT, gap=5, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        p.space_after = Pt(gap)
    return box


def add_round_rect(slide, x, y, w, h, fill=WHITE, line=BORDER, line_w=0.75,
                   radius=0.07, shadow=False):
    shp = slide.shapes.add_shape(RRECT if radius else RECT,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        add_shadow(shp)
    return shp


def add_bar(slide, x, y, w, h, color=BLUE):
    bar = slide.shapes.add_shape(RECT, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_kicker(slide, text, x, y):
    w = 0.34 + len(text) * 0.083
    pill = add_round_rect(slide, x, y, w, 0.34, BLUE, None, radius=0.5)
    add_text(slide, text, x, y + 0.02, w, 0.3, 9, WHITE, True,
             PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return pill


def add_header(slide, kicker, title, page):
    add_kicker(slide, kicker, 0.62, 0.5)
    add_bar(slide, 0.62, 1.04, 0.1, 0.52, BLUE)
    add_text(slide, title, 0.86, 1.04, 11.6, 0.62, 21, SLATE, True, line_spacing=1.0)
    add_text(slide, page, 12.5, 0.54, 0.4, 0.22, 11, SLATE_LIGHT, False, PP_ALIGN.RIGHT)
    add_rule(slide, 0.62, 1.74, 12.1, BORDER)


def add_footer(slide, text="Interview deck · editable PPTX"):
    add_rule(slide, 0.62, 7.08, 12.1, RGBColor(0xE2, 0xE8, 0xF0), 0.8)
    add_text(slide, text, 0.62, 7.17, 7.0, 0.18, 7, SLATE_LIGHT)


def add_rule(slide, x, y, w, color=BORDER, width=1.0):
    line = add_bar(slide, x, y, w, Pt(width) / Inches(1), color)
    return line


def make_qr(label, url):
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    qr = qrcode.QRCode(border=1, box_size=12)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="#16323F", back_color="white")
    path = ASSETS_DIR / f"qr-{label}.png"
    img.save(path)
    return path


def place_qr(slide, label, title, url, x, y, size=1.18, caption_color=SLATE, show_url=True):
    path = make_qr(label, url)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(size), Inches(size))
    add_text(slide, title, x - 0.4, y + size + 0.14, size + 0.8, 0.24, 10, caption_color,
             True, PP_ALIGN.CENTER)
    if show_url:
        add_text(slide, url, x - 0.6, y + size + 0.42, size + 1.2, 0.24, 7, SLATE_LIGHT,
                 False, PP_ALIGN.CENTER)
    return path


def set_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs, bg=BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, bg)
    return slide


# ── slides ───────────────────────────────────────────────────────────
def slide_1(prs):
    slide = blank_slide(prs, WHITE)
    # right blue panel
    panel = add_bar(slide, 9.05, 0, 4.29, 7.5, BLUE)
    add_bar(slide, 9.05, 0, 0.14, 7.5, BLUE_DARK)
    # left content
    add_kicker(slide, "HR OPERATIONS", 0.72, 0.72)
    add_text(slide, "문제를 구조로 바꾸고,\n운영이 흔들리지 않게 만드는\nHR Generalist",
             0.72, 1.62, 8.0, 2.5, 33, SLATE, True, line_spacing=1.04)
    add_bar(slide, 0.74, 4.46, 0.78, 0.06, BLUE)
    add_text(slide, "채용·근태·총무·컴플라이언스·자동화를 연결해\n안정적인 운영을 만듭니다.",
             0.74, 4.66, 7.7, 0.9, 15, SLATE_LIGHT, line_spacing=1.25)
    namebox = add_round_rect(slide, 0.74, 5.95, 3.7, 0.64, WHITE, BORDER, radius=0.5, shadow=True)
    add_text(slide, "성명재   |   HR Operations", 0.74, 6.13, 3.7, 0.28, 13, SLATE, True,
             PP_ALIGN.CENTER)
    # panel content: metrics + QR
    add_text(slide, "핵심 운영 성과", 9.55, 0.78, 3.5, 0.3, 12, SKY, True)
    metrics = [("−85%", "수기 정정 오류"), ("30→111명", "생산직 지원자"), ("0건", "52시간 위반 클레임")]
    for i, (v, l) in enumerate(metrics):
        yy = 1.36 + i * 1.02
        add_text(slide, v, 9.55, yy, 3.4, 0.52, 31, WHITE, True)
        add_text(slide, l, 9.57, yy + 0.56, 3.5, 0.24, 11, SKY)
    qrbox = add_round_rect(slide, 9.55, 4.95, 1.9, 1.95, WHITE, None, radius=0.1, shadow=True)
    place_qr(slide, "evidence", "공개 AI 증거 페이지", QR_URLS["evidence"], 9.92, 5.18, 1.16, show_url=False)
    add_footer(slide, "Conversation deck · evidence-linked")


def slide_2(prs):
    slide = blank_slide(prs)
    add_header(slide, "WHY THIS MATERIAL",
               "오늘은 자동화 자랑이 아니라, 운영의 빈틈을 어떻게 메우는지 보여드립니다.", "02")
    cards = [
        ("회사 요청 ①", "자동화·효율화 사례", "반복 행정과 휴먼에러를 줄인\n운영 구조"),
        ("회사 요청 ②", "강조하고 싶은 기획 경험", "없는 절차를 세우고\n사람이 정착하도록 만든 경험"),
    ]
    for idx, (kicker, title, body) in enumerate(cards):
        x = 1.0 + idx * 5.85
        card = add_round_rect(slide, x, 2.25, 5.25, 2.45, WHITE, BORDER, shadow=True)
        add_bar(slide, x, 2.25, 0.12, 2.45, BLUE)
        add_text(slide, kicker, x + 0.4, 2.55, 2.0, 0.25, 10, BLUE_DARK, True)
        add_text(slide, title, x + 0.4, 2.98, 4.6, 0.4, 21, SLATE, True)
        add_text(slide, body, x + 0.4, 3.62, 4.55, 0.7, 14, SLATE_LIGHT, line_spacing=1.2)
    banner = add_round_rect(slide, 2.1, 5.4, 9.1, 0.82, TINT, TINT_BORDER, radius=0.5)
    add_text(slide, "공통 프레임:  문제  →  구조  →  결과", 2.1, 5.6, 9.1, 0.3, 19, BLUE_DARK,
             True, PP_ALIGN.CENTER)
    add_footer(slide)


def slide_3(prs):
    slide = blank_slide(prs)
    add_header(slide, "OPERATING RANGE",
               "흩어진 일을 한 사람이 연결하면, 회사가 흔들리지 않습니다.", "03")
    lanes = ["채용", "근태·근로시간", "도급·파견", "급여·지원금",
             "노무·감사\n컴플라이언스", "총무·이벤트", "자동화\n(운영 인프라)"]
    start_x, y, w, gap = 0.82, 2.45, 1.6, 0.18
    for idx, lane in enumerate(lanes):
        x = start_x + idx * (w + gap)
        last = idx == 6
        card = add_round_rect(slide, x, y, w, 1.4, BLUE if last else WHITE,
                              BLUE if last else BORDER, shadow=True)
        add_text(slide, lane, x + 0.1, y, w - 0.2, 1.4, 12, WHITE if last else SLATE,
                 True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        if not last:
            add_text(slide, "→", x + w - 0.04, y + 0.5, 0.26, 0.3, 15, BLUE_DARK, True,
                     PP_ALIGN.CENTER)
    band = add_round_rect(slide, 1.4, 4.85, 10.55, 0.95, TINT, TINT_BORDER)
    add_text(slide, "자동화는 별도 영역이 아니라, 앞의 업무들이 흔들리지 않게 받치는 운영 인프라입니다.",
             1.7, 5.12, 9.95, 0.4, 16, BLUE_DARK, True, PP_ALIGN.CENTER)
    add_footer(slide)


def slide_4(prs):
    slide = blank_slide(prs)
    add_header(slide, "MEASURED OUTCOMES", "도구 화면이 아니라, 운영 결과입니다.", "04")
    metrics = [
        ("수기 정정 오류", "−85%", "월 100건 → 15건  (3개월 평균)", False),
        ("52시간 위반 클레임", "0건", "초과예상 144명 → 21명\n특별연장근로 2건 → 0건", True),
        ("생산직 지원자", "30→111명", "공고당 평균 지원자 수", False),
        ("태깅 인식률", "+12%p", "73% → 85%", False),
    ]
    positions = [(0.82, 2.05), (6.86, 2.05), (0.82, 4.5), (6.86, 4.5)]
    for (label, value, desc, green), (x, y) in zip(metrics, positions):
        add_round_rect(slide, x, y, 5.62, 1.96, WHITE, BORDER, shadow=True)
        add_bar(slide, x, y, 0.13, 1.96, GREEN if green else BLUE)
        add_text(slide, label, x + 0.42, y + 0.3, 4.0, 0.25, 11, SLATE_LIGHT, True)
        add_text(slide, value, x + 0.42, y + 0.64, 3.0, 0.7, 40, GREEN if green else BLUE, True)
        add_text(slide, desc, x + 0.42, y + 1.42, 5.0, 0.45, 12.5, SLATE, line_spacing=1.1)
    add_text(slide, "보조 지표 · ATS 전체 지원자 +27%", 0.86, 6.62, 5.0, 0.22, 10, SLATE_LIGHT)
    add_footer(slide)


def slide_5(prs):
    slide = blank_slide(prs)
    add_header(slide, "AUTOMATION METHOD",
               "AI가 한 게 아니라, 제가 질문을 정의하고 검증하고 취사선택했습니다.", "05")
    cases = [
        ("경조화환 자동화", "반복 주문·누락", "신청·확인·발주 흐름 고정", "휴먼에러 0 · 행정 90%↓"),
        ("전자서명 수집", "전사 700여명 회수 지연", "대상·상태·리마인드 체계화", "2주 → 3일"),
        ("직무 키워드 분석", "정성 판단이 흔들림", "페르소나·프레임워크 고정", "외부 600건+ 정량화"),
        ("채용 포스터 생성기", "공고별 포스터 반복", "시트·AI 데이터로 생성", "포스터·카피 직접 생성"),
    ]
    steps = ["문제", "구조", "결과"]
    for idx, (title, prob, struct, result) in enumerate(cases):
        x = 0.78 + (idx % 2) * 6.04
        y = 2.0 + (idx // 2) * 2.18
        add_round_rect(slide, x, y, 5.6, 1.74, WHITE, BORDER, shadow=True)
        add_bar(slide, x, y, 0.12, 1.74, BLUE)
        add_text(slide, title, x + 0.38, y + 0.24, 5.0, 0.26, 14.5, SLATE, True)
        vals = [prob, struct, result]
        colors = [SLATE_LIGHT, SLATE_LIGHT, BLUE_DARK]
        for s in range(3):
            sx = x + 0.4 + s * 1.74
            add_text(slide, steps[s], sx, y + 0.74, 1.6, 0.2, 8, BLUE_DARK, True)
            add_text(slide, vals[s], sx, y + 0.99, 1.6, 0.5, 10.5, colors[s], s == 2,
                     line_spacing=1.05)
            if s < 2:
                add_text(slide, "→", sx + 1.5, y + 0.95, 0.22, 0.24, 12, BLUE, True)
    add_footer(slide)


def slide_6(prs):
    slide = blank_slide(prs)
    add_header(slide, "PLANNING RANGE",
               "큰 조직만 아는 게 아니라, 없는 걸 세우고 사람을 정착시켜본 사람입니다.", "06")
    # card 1
    add_round_rect(slide, 0.84, 2.08, 3.7, 3.05, WHITE, BORDER, shadow=True)
    add_bar(slide, 0.84, 2.08, 0.12, 3.05, BLUE)
    add_text(slide, "공군 학군단 창설 TF", 1.18, 2.4, 3.2, 0.3, 16, SLATE, True)
    add_text(slide, "0 → 1", 1.18, 2.82, 3.2, 0.5, 26, BLUE, True)
    add_multiline(slide, ["조직을 처음부터 구축", "후보생 지원 100%↑", "창설 검열 미흡사항 없음"],
                  1.18, 3.5, 3.2, 1.4, 12.5, SLATE_LIGHT, gap=6)
    # card 2
    add_round_rect(slide, 4.82, 2.08, 3.7, 3.05, WHITE, BORDER, shadow=True)
    add_bar(slide, 4.82, 2.08, 0.12, 3.05, BLUE)
    add_text(slide, "온보딩 개선", 5.16, 2.4, 3.2, 0.3, 16, SLATE, True)
    add_text(slide, "−10%p", 5.16, 2.82, 3.2, 0.5, 26, BLUE, True)
    add_multiline(slide, ["수습 조기퇴사율 감소", "3개월 적응도↑", "들어온 사람이 떠나지 않는 구조"],
                  5.16, 3.5, 3.2, 1.4, 12.5, SLATE_LIGHT, gap=6)
    # card 3 (highlight)
    add_round_rect(slide, 8.8, 2.08, 3.66, 3.05, TINT, TINT_BORDER, shadow=True)
    add_bar(slide, 8.8, 2.08, 0.12, 3.05, BLUE_DARK)
    add_text(slide, "Paytalab 적합성", 9.14, 2.4, 3.1, 0.3, 16, BLUE_DARK, True)
    add_text(slide, "성장기에는 전문가들이 깔끔한 핸드오프와 안정적 운영 리듬을 갖게 하는 역할이 핵심입니다.",
             9.14, 2.92, 3.1, 1.2, 13, SLATE, line_spacing=1.25)
    add_text(slide, "제 역할은 그 빈틈을 메우는 것입니다.", 9.14, 4.4, 3.1, 0.5, 14, BLUE_DARK,
             True, line_spacing=1.15)
    add_footer(slide)


def slide_7(prs):
    slide = blank_slide(prs)
    add_header(slide, "CONVERSATION MAP", "더 궁금하신 부분은 여기서 직접 확인하세요.", "07")
    cards = [
        ("evidence", "공개 AI 증거 페이지", QR_URLS["evidence"], "운영 자동화와 증거 화면", 1.05),
        ("portfolio", "HR 커리어 포트폴리오", QR_URLS["portfolio"], "채용·근태·총무·기획 사례", 7.08),
    ]
    for label, title, url, caption, x in cards:
        add_round_rect(slide, x, 2.1, 5.2, 3.75, WHITE, BORDER, shadow=True)
        place_qr(slide, label, title, url, x + 1.78, 2.5, 1.62)
        add_text(slide, caption, x, 5.3, 5.2, 0.26, 13, SLATE_LIGHT, False, PP_ALIGN.CENTER)
    add_footer(slide, "QR codes are deterministic outputs from the exact URLs")


def build_deck():
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    prs = new_deck()
    for fn in [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7]:
        fn(prs)
    prs.save(PPTX_PATH)


def verify():
    prs = Presentation(PPTX_PATH)
    checks = {
        "pptx": str(PPTX_PATH),
        "font_used": FONT,
        "slide_count": len(prs.slides),
        "pptx_size_bytes": PPTX_PATH.stat().st_size,
        "qr_files": [p.name for p in ASSETS_DIR.glob("qr-*.png")],
    }
    VERIFY_PATH.write_text(json.dumps(checks, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(checks, ensure_ascii=False, indent=2))
    assert checks["slide_count"] == 7
    assert checks["pptx_size_bytes"] > 0


if __name__ == "__main__":
    build_deck()
    verify()
